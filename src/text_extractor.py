import requests
from bs4 import BeautifulSoup
import PyPDF2
from pptx import Presentation
import io

def extract_text_from_url(url):
    """Extracts text content from a URL."""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        
        # Remove script and style elements
        for script_or_style in soup(["script", "style"]):
            script_or_style.decompose()
        
        # Get text
        text = soup.get_text()
        
        # Break into lines and remove leading/trailing space on each
        lines = (line.strip() for line in text.splitlines())
        # Break multi-headlines into a line each
        chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
        # Drop blank lines
        text = '\n'.join(chunk for chunk in chunks if chunk)
        return text
    except requests.exceptions.RequestException as e:
        return f"Error fetching URL: {e}"
    except Exception as e:
        return f"Error parsing URL content: {e}"

def extract_text_from_pdf(file_obj):
    """Extracts text from an uploaded PDF file object."""
    try:
        pdf_reader = PyPDF2.PdfReader(file_obj)
        text = ""
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            text += page.extract_text()
        return text
    except Exception as e:
        return f"Error reading PDF: {e}"

def extract_text_from_pptx(file_obj):
    """Extracts text from an uploaded PPTX file object."""
    try:
        prs = Presentation(file_obj)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"Error reading PPTX: {e}"